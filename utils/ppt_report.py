from io import BytesIO
from datetime import datetime
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


BLUE = RGBColor(30, 95, 216)
DARK = RGBColor(17, 24, 39)
GREY = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)


def rupiah(value):
    return f"Rp {value:,.0f}".replace(",", ".")


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=DARK, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    if align == "center":
        p.alignment = PP_ALIGN.CENTER

    return box


def add_card(slide, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(
        5,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(226, 232, 240)
    return shape


def add_kpi_card(slide, title, value, x, y):
    add_card(slide, x, y, 2.45, 1.25, WHITE)
    add_text(slide, title, x + 0.18, y + 0.18, 2.1, 0.3, 12, False, GREY)
    add_text(slide, value, x + 0.18, y + 0.55, 2.1, 0.45, 20, True, BLUE)


def generate_premium_ppt(
    total_revenue,
    total_transaksi,
    total_item,
    aov,
    top_product,
    top_category,
    top_hour,
    top_payment,
    recommendations,
    chart_images=None,
    logo_path=None
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # SLIDE 1 COVER
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLUE

    add_text(slide, "TOLERANSI KOPI", 0.8, 1.2, 6, 0.5, 24, True, WHITE)
    add_text(slide, "Sales Analytics\nExecutive Report", 0.8, 1.85, 7, 1.4, 44, True, WHITE)
    add_text(slide, "Generated automatically by Toleransi Kopi Analytics Platform", 0.85, 6.55, 7, 0.3, 12, False, WHITE)

    add_card(slide, 8.8, 1.25, 3.5, 3.5, WHITE)

    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            Inches(9.35),
            Inches(1.75),
            Inches(2.4),
            Inches(2.4)
        )

    add_text(slide, "Outlet", 8.9, 5.1, 1, 0.3, 11, False, WHITE)
    add_text(slide, "Toleransi Kopi Palagan", 8.9, 5.45, 3.5, 0.3, 16, True, WHITE)
    add_text(slide, datetime.now().strftime("%d %B %Y"), 8.9, 5.9, 3, 0.3, 11, False, WHITE)

    # SLIDE 2 KPI
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    add_text(slide, "Executive Summary", 0.7, 0.45, 6, 0.5, 34, True, DARK)
    add_text(slide, "Ringkasan performa penjualan utama", 0.75, 1.0, 6, 0.3, 14, False, GREY)

    add_kpi_card(slide, "Total Revenue", rupiah(total_revenue), 0.8, 1.8)
    add_kpi_card(slide, "Total Transaksi", f"{total_transaksi:,.0f}", 3.55, 1.8)
    add_kpi_card(slide, "Item Terjual", f"{total_item:,.0f}", 6.3, 1.8)
    add_kpi_card(slide, "Average Order", rupiah(aov), 9.05, 1.8)

    summary_text = (
        f"Selama periode analisis, Toleransi Kopi mencatat revenue {rupiah(total_revenue)} "
        f"dari {total_transaksi:,.0f} transaksi.\n\n"
        f"Produk terlaris adalah {top_product}, kategori terbaik adalah {top_category}, "
        f"dan jam tersibuk terjadi pada pukul {top_hour}:00."
    )

    add_text(
        slide,
        summary_text,
        1.15, 4.0, 10.4, 1.8,
        17,
        False,
        DARK
    )

    # SLIDE 3 BUSINESS INSIGHT
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    add_text(slide, "AI Business Insight", 0.7, 0.45, 6, 0.5, 34, True, DARK)

    insights = [
        ("Hero Product", f"{top_product} menjadi produk utama untuk strategi promosi."),
        ("Top Category", f"{top_category} menjadi kontributor revenue terbesar."),
        ("Peak Hour", f"Jam tersibuk terjadi pada pukul {top_hour}:00."),
        ("Payment", f"{top_payment} menjadi metode pembayaran dominan.")
    ]

    positions = [(0.8, 1.5), (6.8, 1.5), (0.8, 4.0), (6.8, 4.0)]

    for (title, desc), (x, y) in zip(insights, positions):
        add_card(slide, x, y, 5.3, 1.7, WHITE)
        add_text(slide, title, x + 0.25, y + 0.25, 4.8, 0.3, 18, True, BLUE)
        add_text(slide, desc, x + 0.25, y + 0.75, 4.8, 0.6, 14, False, DARK)

    # SLIDE 4-7 CHARTS
    if chart_images:
        for title, img_buffer in chart_images:
            slide = prs.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = LIGHT_BG

            add_text(slide, title, 0.7, 0.45, 8, 0.5, 30, True, DARK)
            add_text(slide, "Visualisasi data penjualan untuk mendukung pengambilan keputusan bisnis.", 0.75, 0.95, 8, 0.3, 13, False, GREY)

            img_buffer.seek(0)
            slide.shapes.add_picture(
                img_buffer,
                Inches(0.8),
                Inches(1.45),
                Inches(7.7),
                Inches(4.6)
            )

            add_card(slide, 8.9, 1.55, 3.6, 4.4, WHITE)
            add_text(slide, "Key Takeaway", 9.2, 1.85, 3, 0.3, 20, True, BLUE)
            add_text(
                slide,
                "Membantu membaca pola performa\n\n"
                "penjualan\n\n"
                "Menentukan prioritas bisnis\n\n"
                "Mendukung strategi promosi\n\n"
                "Mendukung pengelolaan stok",
                9.2,
                2.35,
                3,
                2.5,
                14,
                False,
                DARK
            )

    # SLIDE RECOMMENDATION
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    add_text(slide, "Strategic Recommendation", 0.7, 0.45, 8, 0.5, 34, True, DARK)
    add_text(slide, "Prioritas tindakan bisnis berdasarkan hasil analisis", 0.75, 1.0, 8, 0.3, 14, False, GREY)

    y = 1.7
    for i, rec in enumerate(recommendations[:6], start=1):
        add_card(slide, 0.85, y, 11.6, 0.65, WHITE)
        add_text(slide, f"{i}", 1.05, y + 0.12, 0.4, 0.25, 14, True, BLUE)
        add_text(slide, rec.replace("**", ""), 1.55, y + 0.12, 10.5, 0.3, 13, False, DARK)
        y += 0.82

    # SLIDE CLOSING
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLUE

    add_text(slide, "THANK YOU", 0, 2.4, 13.3, 0.8, 48, True, WHITE, "center")
    add_text(slide, "Toleransi Kopi Sales Analytics Platform", 0, 3.3, 13.3, 0.4, 18, False, WHITE, "center")
    add_text(slide, "Data-driven insight for better business decisions.", 0, 3.8, 13.3, 0.4, 13, False, WHITE, "center")

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return output